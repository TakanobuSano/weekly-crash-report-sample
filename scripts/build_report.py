#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Crashlytics 週次レポート PowerPoint テンプレート
=================================================
JSON データを受け取り、固定レイアウトで .pptx を生成する。
デザイン・レイアウトはすべてこのスクリプトが持つため、
定期実行タスク側は JSON を作って本スクリプトを呼ぶだけでよい。

使い方:
    python build_crashlytics_report.py --data report_data.json --out report.pptx

依存:
    pip install python-pptx matplotlib
"""

import argparse
import datetime
import json
import math
import re
import sys
import tempfile
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- 設定
# 日本語フォント（環境に合わせてここだけ変更する）
#   Windows: "Yu Gothic UI" / "Meiryo"   macOS: "Hiragino Sans"
FONT_JP = "Yu Gothic UI"

# --- セマンティックカラー（傾向表示用）---
TREND_UP = "CC2936"     # 増加 = 悪化（赤）
TREND_DOWN = "0B7A75"   # 減少 = 改善（ティール）

# --- ブランドカラー ---
# ブランド オレンジ（※仮の値。正式なカラーコードに差し替えるのはこの1行だけでよい）
BRAND_ACCENT = "FF7A00"

BRAND = BRAND_ACCENT    # 表紙アクセント・各シート見出し・最優先カード
BASE_DARK = "232220"       # ダーク基調色: 表紙背景・見出しバー・本文文字色を1色に統一
                           # （実質ニュートラルのチャコール。完全な無彩色はオレンジとの対比で
                           #   青く転んで見えるため、知覚できない約1%だけ暖色側に置いている）
TITLE_BG = BASE_DARK       # 表紙背景（見出しバーと同色）
BODY_BG = "FFF6ED"         # シート2以降の背景（ブランドオレンジの淡色）

# --- 中面のセマンティック配色 ---
KPI_BARS = ["CC2936", "E6820E", "0B7A75", BASE_DARK]    # KPIカード上部バー（4枚目は基調チャコール）
KPI_VALUES = ["CC2936", "E6820E", "0B7A75", BASE_DARK]  # KPI数値の色
TBL_USER_HDR = "E6820E"         # ユーザー影響テーブルのヘッダー
TBL_USER_HDR_TXT = "FFFFFF"
TBL_TECH_HDR = "1B4D4A"         # 技術詳細テーブルのヘッダー
CHART_THIS = "A4232F"           # グラフ: 今週（深めのクリムゾン）
CHART_LAST = "4D5566"           # グラフ: 前週（ダークスレート）
CHART_WBL = "C2C5CC"            # グラフ: 前々週（淡いスレートグレー）
CHART_HBAR = CHART_THIS         # グラフ: 横棒（シート3の「今週」と同色に統一）
CARD_PRI_HDR = BRAND_ACCENT  # 最終シート「今週の最優先」ヘッダー
CARD_PRI_TXT = "FFFFFF"
CARD_MON_HDR = BASE_DARK        # 最終シート「継続監視」ヘッダー（同色に統一）

# --- 背景・カード・補助色 ---
CARD_BG = "FFFFFF"              # カード・ボックス背景
STRIPE = "F9EFE3"               # テーブル偶数行（暖色系）
CARD_LINE = "EFDFCE"            # カードの淡い暖色ボーダー
GRAY_TXT = "5A6275"
DARK_SUB = "D7D5D3"             # 暗い面の上の補助文字色（実質ニュートラル）
DARK_MUTE = "A3A19E"


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.55)            # 左右マージン
CONTENT_W = SLIDE_W - MARGIN * 2
HEADER_H = Inches(0.62)
DATE_RIGHT_MARGIN = Inches(0.45)  # ヘッダー日付の右マージン（0.4in 以上）

ISSUES_PER_DETAIL_SLIDE = 2       # 技術詳細は 1 枚あたり 2 件で固定


# ---------------------------------------------------------------- 汎用
def rgb(hexstr):
    return RGBColor.from_string(hexstr)


def set_run_font(run, size=None, bold=None, color=None, name=FONT_JP, italic=None):
    """run に latin / eastAsia 両方のフォントを設定する。"""
    f = run.font
    if size is not None:
        f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    if italic is not None:
        f.italic = italic
    if color is not None:
        f.color.rgb = rgb(color)
    f.name = name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", name)


def add_rect(slide, x, y, w, h, fill_hex, line_hex=None, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    # テーマスタイル参照(p:style)を除去。残っているとビューアによっては
    # テーマ由来のグラデーションや影が適用され、同色でも見え方が変わるため。
    style = shp._element.find(qn("p:style"))
    if style is not None:
        shp._element.remove(style)
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill_hex)
    if line_hex:
        shp.line.color.rgb = rgb(line_hex)
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = shadow
    return shp


def add_textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, space_after=4, line_spacing=None,
                word_wrap=True):
    """lines = [ [ (text, opts), ... ], ... ]  段落ごとに run のリスト。"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, runs in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space_after is not None:
            p.space_after = Pt(space_after)
        if line_spacing is not None:
            p.line_spacing = line_spacing
        for text, opts in runs:
            r = p.add_run()
            r.text = text
            set_run_font(r, **opts)
    return tb


def trend_text_and_color(pct):
    """前週比(%)から矢印・表示文字列・色を決める。符号と矢印を必ず一致させる。"""
    if pct is None:
        return "—", GRAY_TXT
    if pct > 5:
        arrow, color = "↑", TREND_UP     # 増加=悪化なので赤
    elif pct < -5:
        arrow, color = "↓", TREND_DOWN   # 減少=改善なのでティール
    else:
        arrow, color = "→", GRAY_TXT
    pct_txt = "±0%" if pct == 0 else f"{pct:+d}%"
    return f"{arrow} {pct_txt}", color


def fmt_num(v):
    """1,000 以上の数値はカンマ区切りで表示する。"""
    if isinstance(v, (int, float)) and abs(v) >= 1000:
        return f"{v:,.0f}"
    return str(v)


def estimate_lines(text, width_in, font_pt):
    """セル内テキストの行数を、単語単位の折り返しを再現して概算する。
    全角=半角2文字分。半角英数の連続（クラス名・識別子）は1語として扱い、
    行末に収まらない場合は次行へ送る（実際のワードラップと同じ挙動）。"""
    if not text:
        return 1
    cap = max(2, int(width_in / (font_pt * 0.0075)))
    total = 0
    for seg in str(text).split("\n"):
        # 半角英数記号の連続は折り返し不可の1語、全角・空白は1文字ずつ
        tokens = re.findall(r"[!-~]+|.", seg)
        lines, line = 1, 0
        for tk in tokens:
            w = sum(1 if ord(c) < 0x80 else 2 for c in tk)
            if line + w <= cap:
                line += w
            elif w > cap:
                # 行幅を超える長語はセル幅で強制分割される
                rest = w - (cap - line)
                lines += math.ceil(rest / cap)
                line = rest % cap or cap
            else:
                lines += 1
                line = w
        total += lines
    return total


# ---------------------------------------------------------------- 部品
def content_header(slide, title, period_str, accent_hex):
    """コンテンツスライド共通ヘッダー（背景 + ネイビーバー + 左アクセント + 右に期間）。"""
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BODY_BG)
    add_rect(slide, 0, 0, SLIDE_W, HEADER_H, BASE_DARK)
    add_rect(slide, 0, 0, Inches(0.16), HEADER_H, accent_hex)
    add_textbox(slide, Inches(0.42), Inches(0.06), Inches(8.5), Inches(0.5),
                [[(title, dict(size=22, bold=True, color="FFFFFF"))]],
                anchor=MSO_ANCHOR.MIDDLE)
    date_w = Inches(3.2)
    add_textbox(slide, SLIDE_W - DATE_RIGHT_MARGIN - date_w, Inches(0.06),
                date_w, Inches(0.5),
                [[(period_str, dict(size=12, color=DARK_SUB))]],
                align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def style_table(table, header_hex, col_widths_in, header_pt=12, body_pt=11):
    """ヘッダー色付き白文字 + 偶数行ストライプの共通スタイル。"""
    tbl = table._tbl
    tbl_pr = tbl.find(qn("a:tblPr"))
    if tbl_pr is not None:
        tbl_pr.set("bandRow", "0")
        tbl_pr.set("firstRow", "0")
    for j, w in enumerate(col_widths_in):
        table.columns[j].width = Inches(w)
    for i, row in enumerate(table.rows):
        for cell in row.cells:
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.05)
            cell.fill.solid()
            if i == 0:
                cell.fill.fore_color.rgb = rgb(header_hex)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            else:
                cell.fill.fore_color.rgb = rgb(STRIPE if i % 2 == 0 else "FFFFFF")
                cell.vertical_anchor = MSO_ANCHOR.TOP


def fill_cell(cell, text, size, bold=False, color=None, align=PP_ALIGN.LEFT):
    color = color or BASE_DARK
    tf = cell.text_frame
    tf.word_wrap = True
    segs = str(text).split("\n")
    for i, seg in enumerate(segs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = seg
        set_run_font(r, size=size, bold=bold, color=color)


def chart_bg():
    return "#" + BODY_BG


def insert_chart_image(slide, fig, x, y, w):
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    fig.savefig(tmp.name, dpi=200, bbox_inches="tight",
                facecolor=chart_bg(), edgecolor="none")
    plt.close(fig)
    return slide.shapes.add_picture(tmp.name, x, y, width=w)


# ---------------------------------------------------------------- スライド
def slide_title(prs, data, period_str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, TITLE_BG)
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, BRAND)
    # ISO 週番号（対象期間の終了日基準）を右上に大きく表示
    iso = datetime.date.fromisoformat(data["period"]["end"]).isocalendar()
    iso_w = Inches(3.6)
    add_textbox(slide, SLIDE_W - Inches(0.6) - iso_w, Inches(0.5), iso_w,
                Inches(0.35),
                [[(f"{iso.year} / ISO WEEK", dict(size=13, color=DARK_MUTE))]],
                align=PP_ALIGN.RIGHT)
    add_textbox(slide, SLIDE_W - Inches(0.6) - iso_w, Inches(0.82), iso_w,
                Inches(1.3),
                [[(f"W{iso.week:02d}",
                   dict(size=76, bold=True, color=BRAND))]],
                align=PP_ALIGN.RIGHT)
    x = Inches(0.95)
    add_textbox(slide, x, Inches(2.1), Inches(11.5), Inches(0.9),
                [[("Crashlytics 週次レポート", dict(size=40, bold=True, color="FFFFFF"))]])
    add_textbox(slide, x, Inches(3.05), Inches(11.5), Inches(0.5),
                [[(data["app_name"], dict(size=22, color=DARK_SUB))]])
    add_textbox(slide, x, Inches(3.95), Inches(11.5), Inches(0.9),
                [[(f"対象期間: {period_str}  ({data['period'].get('tz', 'Asia/Tokyo')})",
                   dict(size=15, color="FFFFFF"))],
                 [(f"生成日: {data['generated_at']} (JST)", dict(size=12, color=DARK_SUB))]],
                space_after=8)
    footer = data.get("footer", "iOS 定例会議 / 監視: Firebase Crashlytics")
    add_textbox(slide, x, Inches(6.55), Inches(11.5), Inches(0.4),
                [[(footer, dict(size=11, color=DARK_MUTE))]])


def kpi_card(slide, x, y, w, h, accent, label, scope, value, value_color, sub):
    add_rect(slide, x, y, w, Inches(0.1), accent)
    add_rect(slide, x, y + Inches(0.1), w, h - Inches(0.1), CARD_BG, line_hex=CARD_LINE)
    inner_x = x + Inches(0.16)
    inner_w = w - Inches(0.32)
    add_textbox(slide, inner_x, y + Inches(0.24), inner_w, Inches(0.55),
                [[(label, dict(size=12, bold=True, color=GRAY_TXT))],
                 [(scope, dict(size=9, color=GRAY_TXT))]], space_after=1)
    value = fmt_num(value)
    disp_w = sum(2 if ord(ch) > 0x7F else 1 for ch in str(value))
    value_size = 38 if disp_w <= 7 else (27 if disp_w <= 11 else 20)
    add_textbox(slide, inner_x, y + Inches(0.86), inner_w, Inches(0.8),
                [[(str(value), dict(size=value_size, bold=True, color=value_color))]])
    # 補足行: 10.5pt・標準位置を基本とし、2行までは折り返しで許容する
    # （フォント・開始位置を他カードと揃えるため）。3行以上になる長文
    # （データなし時の理由文など）のみ 9pt に縮小して上詰めで収める。
    disp_sub = sum(2 if ord(ch) > 0x7F else 1 for ch in str(sub))
    capacity_half = Emu(inner_w).inches / (10.5 * 0.00725)  # 10.5pt の半角換算容量
    lines_105 = max(1, -(-disp_sub // max(1, int(capacity_half))))
    sub_long = lines_105 > 2
    add_textbox(slide, inner_x, y + Inches(1.66 if sub_long else 1.78), inner_w,
                Inches(0.52),
                [[(sub, dict(size=9 if sub_long else 10.5, color=BASE_DARK))]],
                space_after=0, line_spacing=1.0)


def slide_summary(prs, data, period_str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(slide, "エグゼクティブサマリー", period_str, BRAND)
    kpi = data["kpi"]
    gap = Inches(0.24)
    card_w = Emu(int((CONTENT_W - gap * 3) / 4))
    card_h = Inches(2.25)
    y = Inches(0.92)

    inner_w_in = Emu(card_w).inches - 0.32

    def fits_one_line(text, pt=10.5):
        disp = sum(2 if ord(ch) > 0x7F else 1 for ch in str(text))
        return disp <= inner_w_in / (pt * 0.00725)

    def numeric_sub(d):
        cur, prev = d.get("current"), d.get("previous")
        if isinstance(cur, (int, float)) and isinstance(prev, (int, float)) and prev:
            diff = cur - prev
            pct = diff / prev * 100
            full = f"前週: {fmt_num(prev)}　前週比 {diff:+,.0f} ({pct:+.1f}%)"
            if fits_one_line(full):
                return full
            # 桁数が多く1行に収まらない場合は、フォントを縮小せず短縮表記にする
            return f"前週: {fmt_num(prev)}　前週比 {pct:+.1f}%"
        return d.get("note", "")

    bars, vals = KPI_BARS, KPI_VALUES
    cards = [
        (bars[0], "FATAL クラッシュ件数", kpi["fatal_events"].get("scope", ""),
         kpi["fatal_events"]["current"], vals[0], numeric_sub(kpi["fatal_events"])),
        (bars[1], "影響ユーザー数", kpi["affected_users"].get("scope", ""),
         kpi["affected_users"]["current"], vals[1], numeric_sub(kpi["affected_users"])),
        (bars[2], "Issue 数 (FATAL)", kpi["issue_count"].get("scope", ""),
         kpi["issue_count"]["current"], vals[2], numeric_sub(kpi["issue_count"])),
    ]
    # 4枚目のカード: kpi.fourth（任意の指標。例: Android=ANR件数, iOS=NON-FATAL件数）
    # を優先し、無ければ旧スキーマの crash_free_rate にフォールバック
    fourth = kpi.get("fourth")
    if fourth is not None:
        cur = fourth.get("current")
        if cur is None:
            cards.append((bars[3], fourth.get("label", "—"), fourth.get("scope", ""),
                          "データなし", GRAY_TXT, fourth.get("note", "")))
        else:
            cards.append((bars[3], fourth.get("label", "—"), fourth.get("scope", ""),
                          cur, vals[3], numeric_sub(fourth)))
    else:
        cfr = kpi.get("crash_free_rate", {})
        cfr_val = cfr.get("current")
        if cfr_val is None:
            cards.append((bars[3], "クラッシュフリー率", cfr.get("scope", ""),
                          "データなし", GRAY_TXT, cfr.get("note", "")))
        else:
            cards.append((bars[3], "クラッシュフリー率", cfr.get("scope", ""),
                          f"{cfr_val}%", vals[3], cfr.get("note", "")))
    for i, c in enumerate(cards):
        kpi_card(slide, MARGIN + (card_w + gap) * i, y, card_w, card_h, *c)

    # サマリー文（高さはコンテンツに合わせる）
    lines = data.get("summary_lines", [])[:4]
    box_y = y + card_h + Inches(0.3)
    est_lines = sum(estimate_lines(s, 11.6, 13) for s in lines)
    box_h = Inches(min(3.6, 0.3 + est_lines * 0.30))
    add_rect(slide, MARGIN, box_y, CONTENT_W, box_h, CARD_BG, line_hex=CARD_LINE)
    paragraphs = [[("■ ", dict(size=12, bold=True, color=BASE_DARK)),
                   (s, dict(size=13, color=BASE_DARK))] for s in lines]
    add_textbox(slide, MARGIN + Inches(0.18), box_y + Inches(0.12),
                CONTENT_W - Inches(0.36), box_h - Inches(0.24),
                paragraphs, space_after=6, line_spacing=1.1)


def slide_trend(prs, data, period_str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(slide, "安定性トレンド（上位 FATAL Issue 直近3週の推移）", period_str, BRAND)
    tc = data["trend_chart"]

    fig, ax = plt.subplots(figsize=(12.4, 4.7))
    fig.patch.set_facecolor(chart_bg())
    ax.set_facecolor(chart_bg())
    n = len(tc["labels"])
    idx = range(n)
    bw = 0.38

    # null = 計測外（前週ランク外など）。0件と区別するため高さ0+「—」表示にする
    def series(values):
        nums = [v if isinstance(v, (int, float)) else 0 for v in values]
        lbls = [str(v) if isinstance(v, (int, float)) else "—" for v in values]
        return nums, lbls

    # 前々週は任意。あれば3系列（時間順に淡グレー→スレート→クリムゾン）、無ければ従来2系列
    has_wbl = tc.get("week_before_last") is not None
    lw_vals, lw_lbls = series(tc["last_week"])
    tw_vals, tw_lbls = series(tc["this_week"])
    if has_wbl:
        wbl_vals, wbl_lbls = series(tc["week_before_last"])
        bw = 0.27
        lbl_pt = 8.5
        b0 = ax.bar([i - bw for i in idx], wbl_vals, bw, color="#" + CHART_WBL,
                    edgecolor="white", linewidth=0.6, label="2 weeks ago")
        b1 = ax.bar(list(idx), lw_vals, bw, color="#" + CHART_LAST,
                    edgecolor="white", linewidth=0.6, label="Last week")
        b2 = ax.bar([i + bw for i in idx], tw_vals, bw, color="#" + CHART_THIS,
                    edgecolor="white", linewidth=0.6, label="This week")
        ax.bar_label(b0, labels=wbl_lbls, fontsize=lbl_pt, padding=2)
        all_vals = wbl_vals + lw_vals + tw_vals
        badge_x = [i + bw for i in idx]
    else:
        bw = 0.38
        lbl_pt = 10
        b1 = ax.bar([i - bw / 2 for i in idx], lw_vals, bw, color="#" + CHART_LAST,
                    edgecolor="white", linewidth=0.6, label="Last week")
        b2 = ax.bar([i + bw / 2 for i in idx], tw_vals, bw, color="#" + CHART_THIS,
                    edgecolor="white", linewidth=0.6, label="This week")
        all_vals = lw_vals + tw_vals
        badge_x = [i + bw / 2 for i in idx]
    ax.bar_label(b1, labels=lw_lbls, fontsize=lbl_pt, padding=2)
    ax.bar_label(b2, labels=tw_lbls, fontsize=lbl_pt, padding=2)

    # 新規バッジ: その週で初めて上位入りした Issue に「new!!」を付ける。
    #  - 今週バー: 前週（あれば前々週も）に出ておらず今週出たもの
    #  - 前週バー: 前々週に出ておらず前週に出たもの（3系列のときのみ）
    def new_this_week(i):
        if isinstance(tc["last_week"][i], (int, float)):
            return False
        if has_wbl and isinstance(tc["week_before_last"][i], (int, float)):
            return False
        return True

    def new_last_week(i):
        return (has_wbl
                and isinstance(tc["last_week"][i], (int, float))
                and not isinstance(tc["week_before_last"][i], (int, float)))

    def badge(x, y):
        ax.annotate("new!!", (x, y), textcoords="offset points",
                    xytext=(0, 18), ha="center", fontsize=11,
                    fontweight="bold", color="#" + TREND_UP)

    max_v = max(all_vals) if all_vals else 0
    if any(new_this_week(i) or new_last_week(i) for i in idx):
        ax.set_ylim(0, max_v * 1.22)   # バッジ分の上余白を確保
    for i in idx:
        if new_this_week(i):
            badge(badge_x[i], tw_vals[i])
        if new_last_week(i):
            badge(i, lw_vals[i])
    ax.set_xticks(list(idx))
    ax.set_xticklabels(tc["labels"], rotation=18, ha="right", fontsize=10)
    ax.set_ylabel("FATAL events", fontsize=11)
    ax.set_title("Top FATAL issues - last 3 weeks", fontsize=13)
    ax.legend(fontsize=10, facecolor=chart_bg(), edgecolor="#BBBBBB")
    ax.grid(axis="y", linestyle=":", alpha=0.5)
    ax.spines[["top", "right"]].set_visible(False)
    fig.tight_layout()

    pic = insert_chart_image(slide, fig, Inches(0.85), Inches(0.95), Inches(11.6))
    insight_y = Inches(0.95) + pic.height + Inches(0.18)
    if tc.get("insight"):
        n_lines = estimate_lines(tc["insight"], Emu(CONTENT_W).inches - 0.4, 14)
        box_h = Inches(0.18 + n_lines * 0.27)
        add_rect(slide, MARGIN, insight_y, Inches(0.1), box_h, CHART_THIS)
        add_textbox(slide, MARGIN + Inches(0.24), insight_y, CONTENT_W - Inches(0.3),
                    box_h,
                    [[(tc["insight"], dict(size=14, bold=True, color=BASE_DARK))]],
                    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)


def slide_user_impact(prs, data, period_str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(slide, "FATAL ユーザー影響（ユーザー視点）", period_str, BRAND)
    issues = data["issues"]
    rows = len(issues) + 1
    col_w = [0.55, 5.55, 1.55, 1.45, 1.75, 1.35]
    table_h = Inches(0.5 + 0.46 * len(issues))
    gf = slide.shapes.add_table(rows, 6, MARGIN, Inches(0.95),
                                Inches(sum(col_w)), table_h)
    table = gf.table
    style_table(table, TBL_USER_HDR, col_w)
    headers = ["#", "概要（何をすると落ちるか）", "影響ユーザー数", "傾向", "チケット/担当", "ステータス"]
    table.rows[0].height = Inches(0.5)
    for j, htxt in enumerate(headers):
        fill_cell(table.cell(0, j), htxt, 12, bold=True,
                  color=TBL_USER_HDR_TXT, align=PP_ALIGN.CENTER)
    for i, iss in enumerate(issues, start=1):
        table.rows[i].height = Inches(0.46)
        if iss.get("new"):
            # 前週ランク外から新たに上位入りした Issue（生成側が new: true を付与）
            trend, tcolor = "新規", TREND_UP
        else:
            trend, tcolor = trend_text_and_color(iss.get("wow_change_pct"))
        fill_cell(table.cell(i, 0), iss["id"], 11, align=PP_ALIGN.CENTER)
        fill_cell(table.cell(i, 1), iss["user_summary"], 11)
        fill_cell(table.cell(i, 2), fmt_num(iss["affected_users"]), 11,
                  align=PP_ALIGN.CENTER)
        fill_cell(table.cell(i, 3), trend, 11, bold=True, color=tcolor,
                  align=PP_ALIGN.CENTER)
        fill_cell(table.cell(i, 4), iss.get("ticket") or "—", 11,
                  align=PP_ALIGN.CENTER)
        fill_cell(table.cell(i, 5), iss.get("status", "OPEN"), 11,
                  align=PP_ALIGN.CENTER)
    add_textbox(slide, MARGIN, Inches(7.05), CONTENT_W, Inches(0.3),
                [[("傾向: 前週比の増減 ±5% 未満は →（横ばい）、前週ランク外からの上位入りは「新規」と表示",
                   dict(size=9.5, color=DARK_MUTE))]],
                space_after=0, line_spacing=1.0)


def hbar_chart(chart):
    fig, ax = plt.subplots(figsize=(6.0, 3.9))
    fig.patch.set_facecolor(chart_bg())
    ax.set_facecolor(chart_bg())
    labels = chart["labels"][::-1]
    values = chart["values"][::-1]
    bars = ax.barh(labels, values, color="#" + CHART_HBAR)
    ax.bar_label(bars, fontsize=10, padding=3)
    ax.set_xlabel("FATAL events", fontsize=11)
    ax.set_title(chart.get("title", ""), fontsize=12)
    ax.tick_params(axis="y", labelsize=10)
    ax.grid(axis="x", linestyle=":", alpha=0.5)
    ax.spines[["top", "right"]].set_visible(False)
    fig.tight_layout()
    return fig


def no_data_panel(slide, x, y, w, h, title, reason):
    add_rect(slide, x, y, w, h, CARD_BG, line_hex=CARD_LINE)
    add_textbox(slide, x + Inches(0.3), y, w - Inches(0.6), h,
                [[(title, dict(size=15, bold=True, color=GRAY_TXT))],
                 [("データなし", dict(size=26, bold=True, color=GRAY_TXT))],
                 [(reason, dict(size=11, color=GRAY_TXT))]],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=8)


def slide_version_os(prs, data, period_str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(slide, "バージョン別・OS 別の安定性比較", period_str, BRAND)
    vc, oc = data.get("version_chart", {}), data.get("os_chart", {})
    half_w = Emu(int((CONTENT_W - Inches(0.3)) / 2))
    y = Inches(1.0)
    panel_h = Inches(4.3)
    notes = []
    for i, (chart, fallback_title) in enumerate(
            [(vc, "バージョン別 (this week)"), (oc, "OS 別 (this week)")]):
        x = MARGIN if i == 0 else MARGIN + half_w + Inches(0.3)
        if chart.get("available"):
            chart.setdefault("title", fallback_title)
            insert_chart_image(slide, hbar_chart(chart), x, y, half_w)
            if chart.get("note"):
                notes.append(chart["note"])
        else:
            no_data_panel(slide, x, y, half_w, panel_h, fallback_title,
                          chart.get("reason", ""))
    if notes:
        add_textbox(slide, MARGIN, Inches(5.65), CONTENT_W, Inches(1.4),
                    [[("※ " + n, dict(size=11, color=GRAY_TXT))] for n in notes],
                    space_after=3)


def slide_tech_details(prs, data, period_str):
    issues = data["issues"]
    col_w = [0.5, 2.6, 3.95, 2.7, 2.45]
    headers = ["#", "該当クラス/メソッド", "スタックトレース要約", "再現条件", "調査メモ"]
    body_pt = 10
    for start in range(0, len(issues), ISSUES_PER_DETAIL_SLIDE):
        chunk = issues[start:start + ISSUES_PER_DETAIL_SLIDE]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        content_header(slide, "FATAL 技術詳細（開発者向け）", period_str, BRAND)
        # 行高さをコンテンツ量から概算（はみ出し防止）
        row_heights = []
        for iss in chunk:
            t = iss["tech"]
            cells = [t["class_method"], t["stack_summary"], t["repro"], t["notes"]]
            lines = max(estimate_lines(c, w - 0.2, body_pt)
                        for c, w in zip(cells, col_w[1:]))
            row_heights.append(min(2.7, 0.25 + lines * 0.185))
        table_h = Inches(0.5 + sum(row_heights))
        gf = slide.shapes.add_table(len(chunk) + 1, 5, MARGIN, Inches(0.95),
                                    Inches(sum(col_w)), table_h)
        table = gf.table
        style_table(table, TBL_TECH_HDR, col_w)
        table.rows[0].height = Inches(0.5)
        for j, htxt in enumerate(headers):
            fill_cell(table.cell(0, j), htxt, 11.5, bold=True, color="FFFFFF",
                      align=PP_ALIGN.CENTER)
        for i, iss in enumerate(chunk, start=1):
            t = iss["tech"]
            table.rows[i].height = Inches(row_heights[i - 1])
            fill_cell(table.cell(i, 0), iss["id"], body_pt, align=PP_ALIGN.CENTER)
            fill_cell(table.cell(i, 1), t["class_method"], body_pt)
            fill_cell(table.cell(i, 2), t["stack_summary"], body_pt)
            fill_cell(table.cell(i, 3), t["repro"], body_pt)
            fill_cell(table.cell(i, 4), t["notes"], body_pt)


def is_light(hexstr):
    r, g, b = (int(hexstr[i:i + 2], 16) for i in (0, 2, 4))
    return (0.299 * r + 0.587 * g + 0.114 * b) > 170


def action_body_height(items, w):
    """アクションカード本文に必要な高さ（Inches）を見積もる。"""
    body_lines = 0
    for it in items:
        body_lines += estimate_lines(it["title"], w.inches - 0.6, 12.5)
        if it.get("detail"):
            body_lines += estimate_lines(it["detail"], w.inches - 0.6, 11.5)
    return Inches(min(4.6, 0.35 + body_lines * 0.34))


def action_card(slide, x, y, w, header_hex, header_txt, title, items,
                body_h=None):
    bullet_hex = BASE_DARK if is_light(header_hex) else header_hex
    add_rect(slide, x, y, w, Inches(0.5), header_hex)
    add_textbox(slide, x + Inches(0.18), y, w - Inches(0.36), Inches(0.5),
                [[(title, dict(size=15, bold=True, color=header_txt))]],
                anchor=MSO_ANCHOR.MIDDLE)
    paragraphs = []
    for it in items:
        ref = "・".join(f"#{n}" for n in it.get("issues", []))
        head = [("● ", dict(size=11, bold=True, color=bullet_hex)),
                (it["title"], dict(size=12.5, bold=True, color=BASE_DARK))]
        if ref:
            head.append((f"（Issue {ref}）", dict(size=11, color=GRAY_TXT)))
        paragraphs.append(head)
        if it.get("detail"):
            paragraphs.append([("　 " + it["detail"], dict(size=11.5, color=GRAY_TXT))])
    if body_h is None:
        body_h = action_body_height(items, w)
    add_rect(slide, x, y + Inches(0.5), w, body_h, CARD_BG, line_hex=CARD_LINE)
    add_textbox(slide, x + Inches(0.22), y + Inches(0.68), w - Inches(0.44),
                body_h - Inches(0.3), paragraphs, space_after=7, line_spacing=1.05)


def slide_actions(prs, data, period_str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(slide, "まとめ & 次のアクション", period_str, BRAND)
    ac = data["actions"]
    half_w = Emu(int((CONTENT_W - Inches(0.3)) / 2))
    y = Inches(1.0)
    pri = ac.get("top_priority", [])
    mon = ac.get("monitoring", [])
    # 2枚のカードは大きい方の高さに揃える
    body_h = Emu(max(action_body_height(pri, half_w),
                     action_body_height(mon, half_w)))
    action_card(slide, MARGIN, y, half_w, CARD_PRI_HDR,
                CARD_PRI_TXT, "今週の最優先", pri, body_h=body_h)
    action_card(slide, MARGIN + half_w + Inches(0.3), y, half_w,
                CARD_MON_HDR, "FFFFFF", "継続監視", mon, body_h=body_h)
    if ac.get("footnote"):
        n = estimate_lines(ac["footnote"], Emu(CONTENT_W).inches, 9.5)
        add_textbox(slide, MARGIN, Inches(7.28) - Inches(0.17) * n, CONTENT_W,
                    Inches(0.17) * n + Inches(0.1),
                    [[(ac["footnote"], dict(size=9.5, color=DARK_MUTE))]],
                    space_after=0, line_spacing=1.0)


# ---------------------------------------------------------------- 検品
def self_check(path):
    """保存後の機械検品: テーブル・図形がスライド下端をはみ出していないか。"""
    prs = Presentation(path)
    problems = []
    for si, slide in enumerate(prs.slides, start=1):
        for shp in slide.shapes:
            if shp.top is None or shp.height is None:
                continue
            bottom = shp.top + shp.height
            if bottom > prs.slide_height + Emu(1):
                problems.append(
                    f"スライド{si}: shape '{shp.shape_type}' が下端を "
                    f"{Emu(bottom - prs.slide_height).inches:.2f}in 超過")
    return problems


# ---------------------------------------------------------------- main
def build(data, out_path):
    period_str = f"{data['period']['start']} 〜 {data['period']['end']}"
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs, data, period_str)
    slide_summary(prs, data, period_str)
    slide_trend(prs, data, period_str)
    slide_user_impact(prs, data, period_str)
    slide_version_os(prs, data, period_str)
    slide_tech_details(prs, data, period_str)
    slide_actions(prs, data, period_str)

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)

    problems = self_check(out_path)
    if problems:
        print("[WARN] レイアウト検品で問題を検出:", file=sys.stderr)
        for p in problems:
            print("  -", p, file=sys.stderr)
        return 2
    print(f"OK: {out_path}（{len(prs.slides.__iter__.__self__._sldIdLst)}枚）")
    return 0


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--data", required=True, help="レポートデータ JSON")
    ap.add_argument("--out", required=True, help="出力 .pptx パス")
    args = ap.parse_args()
    with open(args.data, encoding="utf-8") as f:
        data = json.load(f)
    sys.exit(build(data, args.out))


if __name__ == "__main__":
    main()
